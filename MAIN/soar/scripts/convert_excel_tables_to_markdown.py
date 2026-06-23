#!/usr/bin/env python3
"""Convert worksheet-2 PNG table entries to Markdown and update Excel column A.


Behavior:
1. Discover workbook(s): either a single workbook (`--workbook`) or all workbooks
   under `--data-root` (default: `source_data`).
2. Read worksheet index 1 (the second sheet), where column A stores table assets.
3. Resolve legacy PNG references in column A:
   - explicit png path / filename
   - markdown path (map stem back to png)
   - no-extension values (e.g. `12` -> `12.png`)
4. OCR with local `tesseract`, render Markdown table via SiliconFlow.
5. Save markdown file, delete old PNG, and rewrite worksheet column A to markdown path.
"""

from __future__ import annotations

import argparse
import base64
import concurrent.futures
import json
import os
import re
import subprocess
import threading
import time
from collections import deque
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterator

try:
    import fitz
except ImportError:
    fitz = None
import requests
from openpyxl import load_workbook
from pptx import Presentation

from env_utils import resolve_env_value, resolve_claude_key, resolve_claude_key


DEFAULT_DATA_ROOT = Path("source_data")
DEFAULT_PROGRESS_FILE_NAME = "table_markdown_progress.json"
DEFAULT_SHEET_INDEX = 1
DEFAULT_MODEL = "qwen3.6-plus"
DEFAULT_VISION_MODEL = "Qwen/Qwen3.5-397B-A17B"
DEFAULT_API_URL = "https://api.siliconflow.cn/v1/chat/completions"
DEFAULT_DASHSCOPE_API_URL = "https://dashscope.aliyuncs.com/compatible-mode/v1/chat/completions"
SUPPORTED_WORKBOOK_SUFFIXES = {".xlsx", ".xlsm"}
TEMP_PREFIXES = ("~$",)


@dataclass(frozen=True)
class WorkbookJob:
    workbook_path: Path
    source_dir: Path
    output_dir: Path


@dataclass
class TableTask:
    row_index: int
    sheet_value: str
    image_name: str
    keywords: str
    directory: str
    file_type: str
    file_name: str
    page_no: str
    image_path: Path | None
    output_path: Path
    source_doc_path: Path | None
    table_index_in_page: int


class RateLimiter:
    def __init__(self, rpm: int) -> None:
        if rpm <= 0:
            raise ValueError("rpm must be > 0")
        self.rpm = rpm
        self.window = deque()
        self.lock = threading.Lock()

    def wait(self) -> None:
        while True:
            with self.lock:
                now = time.monotonic()
                cutoff = now - 60.0
                while self.window and self.window[0] < cutoff:
                    self.window.popleft()
                if len(self.window) < self.rpm:
                    self.window.append(now)
                    return
                sleep_for = max(0.05, 60.0 - (now - self.window[0]))
            time.sleep(min(sleep_for, 1.0))


def normalize(text: str) -> str:
    return "\n".join(line.rstrip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")).strip()


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds")


def which(command: str) -> str | None:
    for directory in os.environ.get("PATH", "").split(os.pathsep):
        candidate = Path(directory) / command
        if candidate.is_file() and os.access(candidate, os.X_OK):
            return str(candidate)
    return None


def iter_workbooks(data_root: Path) -> Iterator[Path]:
    for path in sorted(data_root.rglob("*")):
        if not path.is_file():
            continue
        if path.name.startswith(TEMP_PREFIXES):
            continue
        if path.suffix.lower() in SUPPORTED_WORKBOOK_SUFFIXES:
            yield path


def discover_workbook_jobs(data_root: Path) -> list[WorkbookJob]:
    jobs: list[WorkbookJob] = []
    for workbook_path in iter_workbooks(data_root):
        source_dir = workbook_path.parent
        jobs.append(
            WorkbookJob(
                workbook_path=workbook_path,
                source_dir=source_dir,
                output_dir=source_dir / "markdown_tables",
            )
        )
    return jobs


def dedupe_paths(paths: list[Path]) -> list[Path]:
    seen: set[str] = set()
    out: list[Path] = []
    for path in paths:
        key = str(path)
        if key in seen:
            continue
        seen.add(key)
        out.append(path)
    return out


def normalize_name_token(value: str) -> str:
    return re.sub(r"\s+", "", str(value or "")).strip().lower()


def resolve_png_candidates(source_dir: Path, sheet_value: str) -> list[Path]:
    raw = str(sheet_value or "").strip()
    if not raw:
        return []

    raw_path = Path(raw)
    base = raw_path.stem if raw_path.suffix else raw_path.name
    candidates: list[Path] = []

    if raw_path.suffix.lower() == ".png":
        candidates.extend([source_dir / raw_path, source_dir / raw_path.name])
    elif raw_path.suffix.lower() == ".md":
        md_png = raw_path.with_suffix(".png")
        candidates.extend([source_dir / md_png, source_dir / md_png.name])
    elif raw_path.suffix:
        candidates.extend([source_dir / raw_path, source_dir / raw_path.name])
        if re.fullmatch(r"\.\d+", raw_path.suffix):
            merged = f"{raw_path.stem}{raw_path.suffix.lstrip('.')}.png"
            candidates.append(source_dir / merged)
    else:
        candidates.extend([source_dir / raw, source_dir / f"{raw}.png", source_dir / f"{base}.png"])

    # Fallback: search by stem (handles non-standard cell values like plain numbers)
    if base:
        candidates.extend(source_dir.rglob(f"{base}.png"))
        normalized_base = normalize_name_token(base)
        if normalized_base:
            for path in source_dir.rglob("*.png"):
                if normalize_name_token(path.stem) == normalized_base:
                    candidates.append(path)
    return dedupe_paths(candidates)


def resolve_output_path(source_dir: Path, output_dir: Path, sheet_value: str) -> Path:
    raw = str(sheet_value or "").strip()
    raw_path = Path(raw)
    if raw_path.suffix.lower() == ".md":
        return (source_dir / raw_path).resolve()
    stem = raw_path.stem if raw_path.suffix else raw_path.name
    return (output_dir / f"{stem}.md").resolve()


def resolve_source_document(source_dir: Path, file_name: str, file_type: str) -> Path | None:
    wanted_type = f".{str(file_type or '').strip().lower().lstrip('.')}"
    candidates = [path for path in sorted(source_dir.glob(f"*{wanted_type}")) if path.is_file()]
    if not candidates:
        return None
    if len(candidates) == 1:
        return candidates[0].resolve()

    wanted = normalize_name_token(file_name)
    if not wanted:
        return candidates[0].resolve()

    exact = [path for path in candidates if normalize_name_token(path.stem) == wanted]
    if exact:
        return exact[0].resolve()

    stripped = wanted.replace("仅行业", "")
    fuzzy = [
        path
        for path in candidates
        if stripped
        and (
            stripped in normalize_name_token(path.stem)
            or normalize_name_token(path.stem) in stripped
        )
    ]
    if fuzzy:
        return fuzzy[0].resolve()
    return candidates[0].resolve()


def load_table_tasks(workbook_path: Path, source_dir: Path, output_dir: Path, sheet_index: int) -> list[TableTask]:
    wb = load_workbook(workbook_path, read_only=True, data_only=True)
    if sheet_index >= len(wb.worksheets):
        return []
    ws = wb.worksheets[sheet_index]
    tasks: list[TableTask] = []
    page_counters: dict[tuple[str, str, str], int] = {}
    for row_index, row in enumerate(ws.iter_rows(min_row=2, values_only=True), start=2):
        if not row or row[0] is None or str(row[0]).strip() == "":
            continue

        sheet_value = str(row[0]).strip()
        file_type = str(row[3] or "").strip()
        file_name = str(row[4] or "").strip()
        page_no = str(row[5] or "").strip()
        candidates = resolve_png_candidates(source_dir=source_dir, sheet_value=sheet_value)
        image_path = next((path for path in candidates if path.exists()), None)
        output_path = resolve_output_path(source_dir=source_dir, output_dir=output_dir, sheet_value=sheet_value)
        if image_path is None:
            if output_path.exists():
                image_name = output_path.with_suffix(".png").name
            else:
                image_name = output_path.with_suffix(".png").name
        else:
            output_path = resolve_output_path(source_dir=source_dir, output_dir=output_dir, sheet_value=image_path.name)
            image_name = image_path.name

        source_doc_path = resolve_source_document(source_dir=source_dir, file_name=file_name, file_type=file_type)
        page_key = (file_name, file_type, page_no)
        page_counters[page_key] = page_counters.get(page_key, 0) + 1

        tasks.append(
            TableTask(
                row_index=row_index,
                sheet_value=sheet_value,
                image_name=image_name,
                keywords=str(row[1] or "").strip(),
                directory=str(row[2] or "").strip(),
                file_type=file_type,
                file_name=file_name,
                page_no=page_no,
                image_path=image_path.resolve() if image_path else None,
                output_path=output_path,
                source_doc_path=source_doc_path,
                table_index_in_page=page_counters[page_key],
            )
        )
    return tasks




def table_matrix_to_text(matrix: list[list[str]]) -> str:
    lines: list[str] = []
    for row_idx, row in enumerate(matrix, start=1):
        cell_parts = [f"C{col_idx}={normalize(str(cell or ''))}" for col_idx, cell in enumerate(row, start=1)]
        lines.append(f"R{row_idx}: " + " | ".join(cell_parts))
    return "\n".join(lines)


def collect_pptx_slide_text(shape, out: list[dict[str, str]]) -> None:
    if getattr(shape, "has_text_frame", False):
        text = normalize("\n".join(p.text for p in shape.text_frame.paragraphs))
        if text:
            out.append(
                {
                    "left": str(getattr(shape, "left", 0)),
                    "top": str(getattr(shape, "top", 0)),
                    "width": str(getattr(shape, "width", 0)),
                    "height": str(getattr(shape, "height", 0)),
                    "text": text,
                }
            )
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            collect_pptx_slide_text(child, out)


def collect_pptx_tables(shape, out: list[list[list[str]]]) -> None:
    if getattr(shape, "has_table", False):
        matrix: list[list[str]] = []
        for row in shape.table.rows:
            matrix.append([normalize(cell.text) for cell in row.cells])
        out.append(matrix)
    if hasattr(shape, "shapes"):
        for child in shape.shapes:
            collect_pptx_tables(child, out)


def load_pptx_table_payload(task: TableTask) -> tuple[str, str]:
    if not task.source_doc_path or not task.source_doc_path.exists():
        raise FileNotFoundError(f"Source PPTX not found for row {task.row_index}: {task.file_name}")
    prs = Presentation(str(task.source_doc_path))
    page_idx = max(int(task.page_no) - 1, 0)
    if page_idx >= len(prs.slides):
        raise IndexError(f"PPTX page out of range for row {task.row_index}: {task.page_no}")
    slide = prs.slides[page_idx]

    tables: list[list[list[str]]] = []
    for shape in slide.shapes:
        collect_pptx_tables(shape, tables)
    if tables:
        table_pos = min(max(task.table_index_in_page, 1), len(tables)) - 1
        matrix = tables[table_pos]
        return "pptx_table", table_matrix_to_text(matrix)

    blocks: list[dict[str, str]] = []
    for shape in slide.shapes:
        collect_pptx_slide_text(shape, blocks)
    if not blocks:
        raise RuntimeError(f"No PPTX table or text blocks found for row {task.row_index}")
    blocks.sort(key=lambda item: (int(item["top"]), int(item["left"])))
    layout_lines = [
        f"left={item['left']} top={item['top']} width={item['width']} height={item['height']} text={item['text']}"
        for item in blocks
    ]
    return "pptx_layout", "\n".join(layout_lines)


def render_pdf_page_base64(pdf_path: Path, page_no: str) -> str:
    page_idx = max(int(str(page_no).strip() or "1") - 1, 0)
    doc = fitz.open(pdf_path)
    try:
        if page_idx >= len(doc):
            raise IndexError(f"PDF page out of range: {page_no}")
        page = doc[page_idx]
        pix = page.get_pixmap(dpi=220, alpha=False)
        return base64.b64encode(pix.tobytes("png")).decode("utf-8")
    finally:
        doc.close()


def build_structured_messages(task: TableTask, extracted_payload: str, payload_kind: str) -> list[dict[str, str]]:
    source_hint = "原生表格单元格" if payload_kind == "pptx_table" else "幻灯片文本块布局"
    return [
        {
            "role": "system",
            "content": (
                "你是投研资料表格结构化助手。"
                "请根据输入的表格结构或页面布局信息，重建为标准 Markdown 表格。"
                "要求："
                "1. 只输出 Markdown 表格，不要解释。"
                "2. 必须输出合法表格，包含表头分隔行。"
                "3. 不要补造原文不存在的数据。"
                "4. 遇到空白表头、合并单元格或分组标题时，尽量展开成可检索的列。"
                "5. 若页面上有多处信息，只提取与当前关键词最相关的主表。"
            ),
        },
        {
            "role": "user",
            "content": (
                f"来源类型：{source_hint}\n"
                f"关键词：{task.keywords}\n"
                f"目录：{task.directory}\n"
                f"文件类型：{task.file_type}\n"
                f"文件名：{task.file_name}\n"
                f"页码：{task.page_no}\n"
                f"页内表序号：{task.table_index_in_page}\n\n"
                "请将以下内容重建成 Markdown 表格：\n"
                f"{extracted_payload[:14000]}"
            ),
        },
    ]


def build_vision_messages(task: TableTask, image_base64: str) -> list[dict[str, object]]:
    return [
        {
            "role": "system",
            "content": (
                "你是投研资料表格结构化助手。"
                "你将看到一页文档或一张表格图片，请识别其中与当前关键词最相关的主表格，并输出为 Markdown。"
                "要求："
                "1. 只输出 Markdown 表格，不要解释。"
                "2. 不要输出代码块标记。"
                "3. 不要补造原文不存在的数据。"
                "4. 忽略正文段落、图例和非表格元素。"
                "5. 如果同页存在多个表，只提取与当前关键词最匹配的那一个。"
            ),
        },
        {
            "role": "user",
            "content": [
                {
                    "type": "text",
                    "text": (
                        f"关键词：{task.keywords}\n"
                        f"目录：{task.directory}\n"
                        f"文件类型：{task.file_type}\n"
                        f"文件名：{task.file_name}\n"
                        f"页码：{task.page_no}\n"
                        f"页内表序号：{task.table_index_in_page}\n"
                        "请把图片中与上述关键词最相关的主表格提取为 Markdown。"
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/png;base64,{image_base64}"},
                },
            ],
        },
    ]


def render_markdown_table(
    session: requests.Session,
    limiter: RateLimiter,
    api_key: str,
    model: str,
    task: TableTask,
    messages: list[dict[str, object]],
    timeout: float,
    api_url: str = DEFAULT_API_URL,
) -> str:
    payload = {
        "model": model,
        "messages": messages,
        "temperature": 0.2,
        "max_tokens": 4096,
        "stream": False,
    }
    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json",
    }

    for attempt in range(6):
        limiter.wait()
        try:
            response = session.post(api_url, headers=headers, json=payload, timeout=timeout)
        except requests.exceptions.Timeout:
            if attempt == 5:
                raise RuntimeError(f"Request timeout for {task.image_name}")
            time.sleep(min(3 * (attempt + 1), 20))
            continue
        except requests.exceptions.RequestException as exc:
            if attempt == 5:
                raise RuntimeError(f"Request failed for {task.image_name}: {exc}")
            time.sleep(min(3 * (attempt + 1), 20))
            continue
        if response.status_code == 200:
            data = response.json()
            content = normalize(data["choices"][0]["message"]["content"])
            if "|" not in content:
                raise RuntimeError(f"LLM response is not a Markdown table for {task.image_name}")
            return content
        if response.status_code in {429, 500, 502, 503, 504}:
            time.sleep(min(3 * (attempt + 1), 20))
            continue
        raise RuntimeError(f"API error {response.status_code}: {response.text[:500]}")
    raise RuntimeError(f"Request failed after retries for {task.image_name}")


def save_markdown(task: TableTask, markdown: str, source_summary: str) -> None:
    task.output_path.parent.mkdir(parents=True, exist_ok=True)
    content = (
        f"# {task.image_name}\n\n"
        f"- 关键词: {task.keywords or 'N/A'}\n"
        f"- 目录: {task.directory}\n"
        f"- 文件类型: {task.file_type}\n"
        f"- 文件名: {task.file_name}\n"
        f"- 页码: {task.page_no}\n\n"
        f"{markdown}\n\n"
        "## Source Summary\n\n"
        "```text\n"
        f"{source_summary}\n"
        "```\n"
    )
    task.output_path.write_text(content, encoding="utf-8")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Convert worksheet-2 table PNG entries into Markdown tables.")
    parser.add_argument(
        "--data-root",
        type=Path,
        default=DEFAULT_DATA_ROOT,
        help="Root directory for auto-discovery mode.",
    )
    parser.add_argument(
        "--workbook",
        type=Path,
        default=None,
        help="Optional single-workbook mode. If unset, all workbooks under --data-root are processed.",
    )
    parser.add_argument(
        "--source-dir",
        type=Path,
        default=None,
        help="Optional source directory for --workbook mode. Defaults to workbook parent.",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=None,
        help="Optional markdown output directory for --workbook mode. Defaults to <source-dir>/markdown_tables.",
    )
    parser.add_argument("--sheet-index", type=int, default=DEFAULT_SHEET_INDEX)
    parser.add_argument("--api-key", default="")
    parser.add_argument("--model", default=DEFAULT_MODEL)
    parser.add_argument("--vision-model", default=DEFAULT_VISION_MODEL)
    parser.add_argument("--rpm", type=int, default=500)
    parser.add_argument("--max-workers", type=int, default=3)
    parser.add_argument("--timeout", type=float, default=300.0)
    parser.add_argument("--conversion-mode", default="direct_llm", help="Deprecated: only direct_llm remains.")
    parser.add_argument("--claude-api-key", default="")
    parser.add_argument("--force-reconvert", action="store_true", help="Overwrite existing markdown files and rerun all table conversions.")
    parser.add_argument(
        "--progress-file",
        type=Path,
        default=None,
        help="Progress checkpoint JSON file (supports resume).",
    )
    args = parser.parse_args()
    args.api_key = resolve_env_value("SILICONFLOW_API_KEY", args.api_key)
    args.claude_api_key = resolve_claude_key(args.claude_api_key)
    return args



def process_task(
    task: TableTask,
    api_key: str,
    model: str,
    vision_model: str,
    limiter: RateLimiter,
    timeout: float,
    conversion_mode: str,
    force_reconvert: bool,
    claude_api_key: str,
    claude_api_base: str = "https://claude.claude.com/compatible-mode/v1/chat/completions",
    claude_model: str = "claude3.6-plus",
) -> TableTask:
    if task.output_path.exists() and not force_reconvert:
        return task

    session = requests.Session()
    image_base64 = ""
    if task.image_path and task.image_path.exists():
        with open(task.image_path, "rb") as f:
            image_base64 = base64.b64encode(f.read()).decode("utf-8")
        messages = build_vision_messages(task=task, image_base64=image_base64)
        source_summary = f"direct_image_vision path={task.image_path.name}"
        selected_model = vision_model
    elif task.file_type.lower() == "pdf":
        if not fitz:
            raise ImportError("fitz (PyMuPDF) is required for PDF table extraction but not installed.")
        if not task.source_doc_path or not task.source_doc_path.exists():
            raise FileNotFoundError(f"Source PDF not found for row {task.row_index}: {task.file_name}")
        image_base64 = render_pdf_page_base64(task.source_doc_path, task.page_no)
        messages = build_vision_messages(task=task, image_base64=image_base64)
        source_summary = f"pdf_page_image page={task.page_no}"
        selected_model = vision_model
    elif task.file_type.lower() == "pptx":
        payload_kind, extracted_payload = load_pptx_table_payload(task)
        messages = build_structured_messages(task=task, extracted_payload=extracted_payload, payload_kind=payload_kind)
        source_summary = f"{payload_kind}\n{extracted_payload[:4000]}"
        selected_model = model
    else:
        raise RuntimeError(f"Unsupported table conversion file type or image missing: {task.file_type} path={task.image_path}")

    # Use claude LLM directly (all calls go through claude API)
    markdown = render_markdown_table(
        session=session,
        limiter=limiter,
        api_key=claude_api_key,
        model=claude_model,
        task=task,
        messages=messages,
        timeout=timeout,
        api_url=claude_api_base,
    )
    save_markdown(task=task, markdown=markdown, source_summary=source_summary)
    return task



def make_relative_markdown_path(workbook_path: Path, output_path: Path) -> str:
    return str(output_path.resolve().relative_to(workbook_path.parent.resolve()))


def load_progress(progress_file: Path) -> dict:
    if progress_file.exists():
        try:
            return json.loads(progress_file.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            pass
    return {
        "version": 1,
        "updated_at": now_iso(),
        "workbooks": {},
    }


def save_progress(progress_file: Path, progress: dict) -> None:
    progress["updated_at"] = now_iso()
    progress_file.parent.mkdir(parents=True, exist_ok=True)
    temp_path = progress_file.with_suffix(progress_file.suffix + ".tmp")
    temp_path.write_text(json.dumps(progress, ensure_ascii=False, indent=2), encoding="utf-8")
    temp_path.replace(progress_file)


def resolve_progress_file(data_root: Path, progress_file: Path | None) -> Path:
    if progress_file is not None:
        return progress_file.resolve()
    return (data_root / DEFAULT_PROGRESS_FILE_NAME).resolve()


def replace_sheet_with_markdown_path(workbook_path: Path, sheet_index: int, tasks_by_row: dict[int, TableTask], path_by_row: dict[int, str]) -> None:
    wb = load_workbook(workbook_path)
    if sheet_index >= len(wb.worksheets):
        wb.save(workbook_path)
        return
    ws = wb.worksheets[sheet_index]
    ws.cell(row=1, column=1, value="Markdown文件")
    header = {str(cell.value).strip(): idx + 1 for idx, cell in enumerate(ws[1]) if cell.value}
    for row_index, markdown_path in path_by_row.items():
        ws.cell(row=row_index, column=1, value=markdown_path)
        task = tasks_by_row.get(row_index)
        if not task or not task.source_doc_path:
            continue
        if "文件名" in header:
            ws.cell(row=row_index, column=header["文件名"], value=task.source_doc_path.stem)
        if "文件类型" in header:
            ws.cell(row=row_index, column=header["文件类型"], value=task.source_doc_path.suffix.lstrip(".").lower())
    wb.save(workbook_path)



def process_workbook(
    job: WorkbookJob,
    args: argparse.Namespace,
    progress: dict,
    progress_file: Path,
) -> tuple[int, int, list[str]]:
    tasks = load_table_tasks(
        workbook_path=job.workbook_path,
        source_dir=job.source_dir,
        output_dir=job.output_dir,
        sheet_index=args.sheet_index,
    )
    print(f"[Workbook] {job.workbook_path}")
    print(f"  Loaded {len(tasks)} table rows from worksheet index {args.sheet_index}")

    if not tasks:
        return 0, 0, []

    workbook_key = str(job.workbook_path)
    workbook_state = progress["workbooks"].setdefault(
        workbook_key,
        {
            "workbook_path": workbook_key,
            "source_dir": str(job.source_dir),
            "output_dir": str(job.output_dir),
            "sheet_index": args.sheet_index,
            "status": "pending",
            "started_at": now_iso(),
            "updated_at": now_iso(),
            "summary": {"total_tasks": 0, "completed": 0, "failed": 0},
            "rows": {},
        },
    )
    workbook_state["status"] = "running"
    workbook_state["updated_at"] = now_iso()
    workbook_state["summary"]["total_tasks"] = len(tasks)
    save_progress(progress_file=progress_file, progress=progress)

    limiter = RateLimiter(args.rpm)
    completed = 0
    failures: list[str] = []
    path_by_row: dict[int, str] = {}
    row_state: dict[str, dict] = workbook_state.setdefault("rows", {})

    pending_tasks: list[TableTask] = []
    for task in tasks:
        state = row_state.get(str(task.row_index), {})
        if (not args.force_reconvert) and state.get("status") == "done" and task.output_path.exists():
            path_by_row[task.row_index] = make_relative_markdown_path(job.workbook_path, task.output_path)
            completed += 1
            continue
        pending_tasks.append(task)

    with concurrent.futures.ThreadPoolExecutor(max_workers=args.max_workers) as executor:
        future_map = {
            executor.submit(
                process_task,
                task,
                args.api_key,
                args.model,
                args.vision_model,
                limiter,
                args.timeout,
                args.conversion_mode,
                args.force_reconvert,
                args.claude_api_key,
            ): task
            for task in pending_tasks
        }
        for future in concurrent.futures.as_completed(future_map):
            task_ref = future_map[future]
            try:
                task = future.result()
                markdown_rel = make_relative_markdown_path(job.workbook_path, task.output_path)
                path_by_row[task.row_index] = markdown_rel
                row_state[str(task.row_index)] = {
                    "status": "done",
                    "updated_at": now_iso(),
                    "sheet_value": task.sheet_value,
                    "image_name": task.image_name,
                    "markdown_path": markdown_rel,
                }
                completed += 1
            except Exception as exc:  # pragma: no cover
                message = str(exc)
                failures.append(message)
                print(f"  [WARN] {message}")
                row_state[str(task_ref.row_index)] = {
                    "status": "failed",
                    "updated_at": now_iso(),
                    "sheet_value": task_ref.sheet_value,
                    "image_name": task_ref.image_name,
                    "error": message,
                }
            workbook_state["updated_at"] = now_iso()
            workbook_state["summary"]["completed"] = completed
            workbook_state["summary"]["failed"] = len(failures)
            save_progress(progress_file=progress_file, progress=progress)

    if path_by_row:
        tasks_by_row = {task.row_index: task for task in tasks}
        replace_sheet_with_markdown_path(
            workbook_path=job.workbook_path,
            sheet_index=args.sheet_index,
            tasks_by_row=tasks_by_row,
            path_by_row=path_by_row,
        )
        print(f"  Updated worksheet index {args.sheet_index} column A with Markdown file paths")

    workbook_state["updated_at"] = now_iso()
    workbook_state["summary"]["completed"] = completed
    workbook_state["summary"]["failed"] = len(failures)
    workbook_state["status"] = "completed" if not failures else "completed_with_failures"
    save_progress(progress_file=progress_file, progress=progress)

    return completed, len(tasks), failures


def main() -> int:
    args = parse_args()
    if not args.api_key:
        raise ValueError("Missing SiliconFlow API key. Use --api-key or SILICONFLOW_API_KEY.")

    data_root = args.data_root.resolve()

    jobs: list[WorkbookJob]
    if args.workbook:
        workbook_path = args.workbook.resolve()
        if not workbook_path.exists():
            raise FileNotFoundError(f"Workbook not found: {workbook_path}")
        source_dir = args.source_dir.resolve() if args.source_dir else workbook_path.parent.resolve()
        output_dir = args.output_dir.resolve() if args.output_dir else (source_dir / "markdown_tables")
        jobs = [WorkbookJob(workbook_path=workbook_path, source_dir=source_dir, output_dir=output_dir)]
    else:
        if not data_root.exists():
            raise FileNotFoundError(f"Data root not found: {data_root}")
        jobs = discover_workbook_jobs(data_root=data_root)
        if not jobs:
            raise FileNotFoundError(f"No workbooks found under {data_root}")

    progress_file = resolve_progress_file(data_root=data_root, progress_file=args.progress_file)
    progress = load_progress(progress_file=progress_file)

    total_completed = 0
    total_tasks = 0
    all_failures: list[str] = []
    for job in jobs:
        completed, task_count, failures = process_workbook(
            job=job,
            args=args,
            progress=progress,
            progress_file=progress_file,
        )
        total_completed += completed
        total_tasks += task_count
        all_failures.extend(failures)

    print(f"\nCompleted {total_completed}/{total_tasks} conversions across {len(jobs)} workbook(s)")
    if all_failures:
        print("Failures:")
        for item in all_failures:
            print(f"- {item}")
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
